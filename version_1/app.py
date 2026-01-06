import streamlit as st
from snowflake.snowpark.context import get_active_session
import uuid, re, datetime as dt
import pandas as pd
import io
import time

# -------------  PAGE CONFIG & STYLING -------------
st.set_page_config(
    page_title="AI Document Assistant",
    page_icon="💬",
    layout="wide",
    initial_sidebar_state="expanded",
)

# Dark theme configuration (only theme)
theme = {
    "primary": "#8b5cf6",  # Purple
    "secondary": "#a78bfa",
    "gradient_start": "#8b5cf6",
    "gradient_end": "#7c3aed",
    "bg_main": "#0f172a",  # Slate dark
    "bg_secondary": "#1e293b",
    "bg_sidebar": "#1a1f2e",
    "text_primary": "#f1f5f9",
    "text_secondary": "#94a3b8",
    "border": "#334155",
    "input_bg": "#1e293b",
    "input_border": "#334155",
    "input_focus": "#8b5cf6",
    "message_user_bg": "linear-gradient(135deg, #8b5cf6 0%, #7c3aed 100%)",
    "message_user_text": "#ffffff",
    "message_ai_bg": "#1e293b",
    "message_ai_text": "#f1f5f9",
    "message_ai_border": "#334155",
    "button_hover_shadow": "rgba(139, 92, 246, 0.4)",
    "welcome_text": "#94a3b8",
    "file_upload_bg": "#1e293b",
    "file_upload_border": "#334155",
    "expander_bg": "#1e293b",
    "accent_highlight": "#a78bfa",  # Light purple accent
    "success_color": "#10b981",  # Green
    "card_shadow": "0 4px 6px -1px rgba(139, 92, 246, 0.2), 0 2px 4px -1px rgba(139, 92, 246, 0.1)",
}

# Global styling
st.markdown(
    f"""
    <style>
    /* Hide Streamlit branding */
    #MainMenu {{visibility: hidden;}}
    footer {{visibility: hidden;}}
    
    /* Main app background with subtle gradient */
    .stApp {{
        background: {theme['bg_main']} !important;
    }}
    
    /* Sidebar with elegant design */
    section[data-testid="stSidebar"] {{
        background: linear-gradient(180deg, {theme['bg_sidebar']} 0%, {theme['bg_main']} 100%) !important;
        border-right: 1px solid {theme['border']} !important;
        box-shadow: 2px 0 12px rgba(0,0,0,0.03);
    }}
    
    section[data-testid="stSidebar"] * {{
        color: {theme['text_primary']} !important;
    }}
    
    /* All text colors */
    .stMarkdown, .stText, p, span, div {{
        color: {theme['text_primary']} !important;
    }}
    
    h1, h2, h3, h4, h5, h6 {{
        color: {theme['text_primary']} !important;
        font-weight: 600 !important;
    }}
    
    /* Enhanced button styling with better shadows */
    .stButton > button {{
        border-radius: 14px !important;
        font-weight: 500 !important;
        transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1) !important;
        border: none !important;
        color: {theme['text_primary']} !important;
        background-color: {theme['bg_secondary']} !important;
        box-shadow: 0 1px 3px rgba(0,0,0,0.08) !important;
        padding: 0.6rem 1.2rem !important;
        font-size: 0.95rem !important;
    }}
    
    .stButton > button:hover {{
        transform: translateY(-2px);
        box-shadow: 0 8px 16px {theme['button_hover_shadow']} !important;
        background-color: {theme['expander_bg']} !important;
    }}
    
    /* Primary buttons with gradient and glow */
    .stButton > button[kind="primary"] {{
        background: {theme['message_user_bg']} !important;
        color: white !important;
        font-weight: 600 !important;
        box-shadow: 0 4px 12px {theme['button_hover_shadow']} !important;
    }}
    
    .stButton > button[kind="primary"]:hover {{
        box-shadow: 0 8px 20px {theme['button_hover_shadow']} !important;
        transform: translateY(-3px);
    }}
    
    /* Enhanced input fields with modern styling */
    .stTextInput > div > div > input,
    .stTextArea > div > div > textarea {{
        border-radius: 16px !important;
        border: 2px solid {theme['input_border']} !important;
        font-size: 15px !important;
        background-color: {theme['input_bg']} !important;
        color: {theme['text_primary']} !important;
        padding: 14px 18px !important;
        transition: all 0.3s ease !important;
        box-shadow: 0 2px 4px rgba(0,0,0,0.02) !important;
    }}
    
    .stTextInput > div > div > input:focus,
    .stTextArea > div > div > textarea:focus {{
        border-color: {theme['input_focus']} !important;
        box-shadow: 0 0 0 4px rgba(139, 92, 246, 0.12), 
                   0 4px 8px rgba(0,0,0,0.08) !important;
        outline: none !important;
        transform: translateY(-1px);
    }}
    
    /* Select box with modern design */
    div[data-testid="stSelectbox"] {{
        color: {theme['text_primary']} !important;
    }}
    
    div[data-testid="stSelectbox"] > div > div {{
        border-radius: 14px !important;
        border: 2px solid {theme['input_border']} !important;
        background-color: {theme['input_bg']} !important;
        box-shadow: 0 2px 4px rgba(0,0,0,0.02) !important;
        transition: all 0.3s ease !important;
    }}
    
    div[data-testid="stSelectbox"] > div > div:hover {{
        border-color: {theme['input_focus']} !important;
        box-shadow: 0 4px 8px rgba(0,0,0,0.06) !important;
    }}
    
    /* Enhanced file uploader */
    section[data-testid="stFileUploader"] {{
        border-radius: 16px !important;
        border: 2px dashed {theme['file_upload_border']} !important;
        background: {theme['file_upload_bg']} !important;
        padding: 24px;
        transition: all 0.3s ease !important;
    }}
    
    section[data-testid="stFileUploader"]:hover {{
        border-color: {theme['input_focus']} !important;
        background: {theme['bg_secondary']} !important;
        transform: translateY(-2px);
        box-shadow: 0 8px 16px rgba(0,0,0,0.06) !important;
    }}
    
    /* Beautiful expanders */
    .streamlit-expanderHeader {{
        border-radius: 12px !important;
        background: {theme['expander_bg']} !important;
        font-weight: 500 !important;
        color: {theme['text_primary']} !important;
        border: 1px solid {theme['border']} !important;
        padding: 12px 16px !important;
        transition: all 0.3s ease !important;
    }}
    
    .streamlit-expanderHeader:hover {{
        background: {theme['bg_secondary']} !important;
        border-color: {theme['input_focus']} !important;
        box-shadow: 0 4px 8px rgba(0,0,0,0.05) !important;
    }}
    
    /* Info/Warning boxes with better design */
    .stAlert {{
        background-color: {theme['bg_secondary']} !important;
        color: {theme['text_primary']} !important;
        border: 1px solid {theme['border']} !important;
        border-radius: 12px !important;
        padding: 16px !important;
        box-shadow: 0 2px 8px rgba(0,0,0,0.04) !important;
    }}
    
    /* Captions with better styling */
    .stCaptionContainer {{
        color: {theme['text_secondary']} !important;
        font-size: 0.9rem !important;
    }}
    
    /* Enhanced progress bar */
    .stProgress > div > div > div {{
        background: {theme['message_user_bg']} !important;
        border-radius: 10px !important;
    }}
    
    .stProgress > div > div {{
        background-color: {theme['bg_secondary']} !important;
        border-radius: 10px !important;
    }}
    
    /* Checkbox styling */
    .stCheckbox {{
        color: {theme['text_primary']} !important;
    }}
    
    .stCheckbox > label {{
        font-weight: 500 !important;
    }}
    
    /* Form container */
    .stForm {{
        border: none !important;
        background-color: transparent !important;
    }}
    
    /* Divider styling */
    hr {{
        margin: 2rem 0 !important;
        border: none !important;
        height: 1px !important;
        background: linear-gradient(90deg, transparent, {theme['border']}, transparent) !important;
    }}
    
    /* Improve spacing */
    .block-container {{
        padding-top: 2rem !important;
        max-width: 1200px !important;
        padding-bottom: 3rem !important;
    }}
    
    /* Metric styling */
    [data-testid="stMetric"] {{
        background: {theme['bg_secondary']} !important;
        padding: 16px !important;
        border-radius: 12px !important;
        border: 1px solid {theme['border']} !important;
    }}
    
    /* Tab styling */
    .stTabs [data-baseweb="tab-list"] {{
        gap: 8px;
    }}
    
    .stTabs [data-baseweb="tab"] {{
        border-radius: 10px !important;
        padding: 10px 20px !important;
        font-weight: 500 !important;
    }}
    
    /* Scrollbar styling */
    ::-webkit-scrollbar {{
        width: 10px;
        height: 10px;
    }}
    
    ::-webkit-scrollbar-track {{
        background: {theme['bg_secondary']};
        border-radius: 10px;
    }}
    
    ::-webkit-scrollbar-thumb {{
        background: {theme['border']};
        border-radius: 10px;
    }}
    
    ::-webkit-scrollbar-thumb:hover {{
        background: {theme['input_border']};
    }}
    </style>
""",
    unsafe_allow_html=True,
)

# -------------  SAFE SESSION -------------
try:
    session = get_active_session()
except Exception as e:
    st.error(f"Cannot obtain Snowflake session: {e}")
    st.stop()

# -------------  CONSTANTS -------------
LLM_OPTIONS = ["mistral-7b", "snowflake-arctic", "mixtral-8x7b"]
DOCS_TBL = "DOCUMENTS"
CHUNKS_TBL = "CHUNKS"
EMB_TBL = "EMBEDDINGS"
CHAT_TBL = "CHAT_HISTORY"

# Supported file extensions
SUPPORTED_EXTENSIONS = {
    "pdf": "PDF",
    "pptx": "PowerPoint",
    "csv": "CSV",
    "xlsx": "Excel",
    "xls": "Excel",
}

# Performance tuning
MAX_CHUNK_SIZE = 1000  # Safer size for 512 token limit
CHUNK_OVERLAP = 200  # Adjusted overlap


# -------------  UTILS -------------
def sanitize(txt: str) -> str:
    return txt.replace("'", "''") if txt else ""


# -------------  SESSION INIT -------------
for k in ("sessions", "current", "llm", "last_upload_time", "kb_refresh"):
    st.session_state.setdefault(
        k,
        (
            {}
            if k == "sessions"
            else (
                LLM_OPTIONS[0]
                if k == "llm"
                else (
                    0
                    if k == "last_upload_time"
                    else False if k == "kb_refresh" else None
                )
            )
        ),
    )

# Validate and reset LLM if it's no longer available
if st.session_state.llm not in LLM_OPTIONS:
    st.session_state.llm = LLM_OPTIONS[0]

# Initialize current session if needed
if not st.session_state.current:
    st.session_state.current = str(uuid.uuid4())
    st.session_state.sessions[st.session_state.current] = {
        "title": "New chat",
        "messages": [],
    }


# -------------  FULL EXTRACTION -------------
def extract_pdf_full(f):
    """Full PDF extraction"""
    try:
        from PyPDF2 import PdfReader

        reader = PdfReader(f)
        text_parts = []
        # Read ALL pages
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        return "".join(text_parts)
    except Exception as e:
        st.error(f"PDF extraction error: {e}")
        return ""


def extract_pptx_full(f):
    """Full PowerPoint extraction"""
    try:
        from pptx import Presentation

        prs = Presentation(f)
        text_parts = []
        for slide in prs.slides:
            slide_text = " ".join(
                shape.text
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            )
            if slide_text:
                text_parts.append(slide_text)
        return "\n\n".join(text_parts)
    except Exception as e:
        st.error(f"PowerPoint extraction error: {e}")
        return ""


def extract_csv_full(f):
    """Full CSV extraction (reads as raw text)"""
    try:
        # Reset file pointer and read as raw text
        f.seek(0)
        return f.read().decode("utf-8")
    except Exception as e:
        st.error(f"CSV extraction error: {e}")
        return ""


def extract_excel_full(f):
    """Full Excel extraction (all sheets, all rows)"""
    try:
        excel_file = pd.ExcelFile(f)
        parts = []

        # Process ALL sheets
        for sheet_name in excel_file.sheet_names:
            parts.append(f"=== Sheet: {sheet_name} ===")
            # Read ALL rows
            df = pd.read_excel(excel_file, sheet_name=sheet_name)

            # Convert the entire DataFrame to string *without* truncation
            # This ensures every character from the Excel file is read
            with pd.option_context(
                "display.max_rows",
                None,
                "display.max_columns",
                None,
                "display.width",
                None,
            ):
                parts.append(df.to_string())

            parts.append("=" * 50)

        return "\n".join(parts)
    except Exception as e:
        st.error(f"Excel extraction error: {e}")
        return ""


def extract_file_content(file, file_type):
    """Extract content based on file type - FULL versions"""
    if file_type.upper() == "PDF":
        return extract_pdf_full(file)
    elif file_type.upper() == "POWERPOINT":
        return extract_pptx_full(file)
    elif file_type.upper() == "CSV":
        return extract_csv_full(file)
    elif file_type.upper() == "EXCEL":
        return extract_excel_full(file)
    else:
        return ""


def chunk_text(
    txt: str, size: int = MAX_CHUNK_SIZE, overlap: int = CHUNK_OVERLAP
) -> list[str]:
    """Split text into overlapping chunks - optimized"""
    if not txt:
        return []
    chunks = []
    start = 0
    txt_len = len(txt)

    while start < txt_len:
        end = min(start + size, txt_len)
        chunks.append(txt[start:end])
        if end == txt_len:
            break
        start += size - overlap

    return chunks


def store_document_fast(
    filename: str, src_type: str, content: str, public: bool = False
) -> bool:
    """
    OPTIMIZED: Store document, chunks, and embeddings using server-side processing.
    NOW WITH IMMEDIATE COMMIT for instant knowledge base update.
    """
    doc_id = str(uuid.uuid4())
    try:
        # 1. Insert document (single operation) with immediate commit
        session.sql(
            f"INSERT INTO {DOCS_TBL} (DOC_ID,FILENAME,FILE_TYPE,FILE_SIZE,IS_PUBLIC,UPLOADED_AT) "
            f"VALUES ('{sanitize(doc_id)}','{sanitize(filename)}','{src_type}',{len(content)},"
            f"{public},CURRENT_TIMESTAMP())"
        ).collect()

        # FORCE COMMIT - Critical for immediate visibility
        session.sql("COMMIT").collect()

        # 2. Create and batch insert chunks
        chunks = chunk_text(content)
        if not chunks:
            st.warning(f"⚠️ {filename}: Document stored but no content to chunk")
            return True  # Document stored, but no content to chunk

        # Prepare chunk data
        chunk_data = []
        for idx, chunk in enumerate(chunks):
            if not chunk.strip():
                continue
            chunk_data.append(
                {
                    "CHUNK_ID": f"{doc_id}_{idx}",
                    "DOC_ID": doc_id,
                    "CHUNK_INDEX": idx,
                    "CHUNK_TEXT": chunk.strip(),
                }
            )

        # Batch insert chunks
        if chunk_data:
            chunks_df = pd.DataFrame(chunk_data)
            session.write_pandas(
                chunks_df,
                table_name=CHUNKS_TBL,
                auto_create_table=False,
                overwrite=False,
                quote_identifiers=False,
            )

            # FORCE COMMIT after chunks
            session.sql("COMMIT").collect()

            # 3. Generate and store embeddings
            session.sql(
                f"""
                INSERT INTO {EMB_TBL} (CHUNK_ID, EMBEDDING)
                SELECT
                    c.CHUNK_ID,
                    SNOWFLAKE.CORTEX.EMBED_TEXT_768('snowflake-arctic-embed-m', LEFT(c.CHUNK_TEXT, 10000))
                FROM
                    {CHUNKS_TBL} c
                WHERE
                    c.DOC_ID = '{sanitize(doc_id)}'
            """
            ).collect()

            # FORCE COMMIT after embeddings
            session.sql("COMMIT").collect()

            # Verify the document was stored
            verify = session.sql(
                f"SELECT COUNT(*) as cnt FROM {DOCS_TBL} WHERE DOC_ID='{sanitize(doc_id)}'"
            ).collect()
            if verify and verify[0][0] > 0:
                return True
            else:
                st.error(
                    f"❌ {filename}: Verification failed - document not found in database"
                )
                return False

        return True

    except Exception as e:
        st.error(f"❌ Storage error for {filename}: {str(e)}")
        # Best-effort cleanup on failure
        try:
            session.sql(
                f"DELETE FROM {EMB_TBL} WHERE CHUNK_ID IN (SELECT CHUNK_ID FROM {CHUNKS_TBL} WHERE DOC_ID='{sanitize(doc_id)}')"
            ).collect()
            session.sql(
                f"DELETE FROM {CHUNKS_TBL} WHERE DOC_ID='{sanitize(doc_id)}'"
            ).collect()
            session.sql(
                f"DELETE FROM {DOCS_TBL} WHERE DOC_ID='{sanitize(doc_id)}'"
            ).collect()
            session.sql("COMMIT").collect()  # Commit cleanup
        except Exception as cleanup_e:
            st.error(f"❌ Cleanup failed for {filename}: {str(cleanup_e)}")
        return False


def delete_document(doc_id: str) -> bool:
    """Delete document and related data"""
    try:
        session.sql(
            f"DELETE FROM {EMB_TBL} WHERE CHUNK_ID IN (SELECT CHUNK_ID FROM {CHUNKS_TBL} WHERE DOC_ID='{sanitize(doc_id)}')"
        ).collect()
        session.sql(
            f"DELETE FROM {CHUNKS_TBL} WHERE DOC_ID='{sanitize(doc_id)}'"
        ).collect()
        session.sql(
            f"DELETE FROM {DOCS_TBL} WHERE DOC_ID='{sanitize(doc_id)}'"
        ).collect()
        return True
    except Exception as e:
        st.error(f"Delete error: {e}")
        return False


def get_user_docs() -> pd.DataFrame:
    """Get user documents"""
    try:
        return session.sql(
            f"SELECT DOC_ID, FILENAME, FILE_TYPE, FILE_SIZE, UPLOADED_AT FROM {DOCS_TBL} ORDER BY UPLOADED_AT DESC"
        ).to_pandas()
    except:
        return pd.DataFrame()


def load_chat_history(session_id: str):
    """Load chat history from database"""
    try:
        rows = session.sql(
            f"SELECT QUERY_TEXT, RESPONSE_TEXT FROM {CHAT_TBL} "
            f"WHERE SESSION_ID='{sanitize(session_id)}' ORDER BY QUERY_TIMESTAMP"
        ).to_pandas()

        messages = []
        for _, r in rows.iterrows():
            messages.append({"role": "user", "content": r["QUERY_TEXT"]})
            messages.append({"role": "assistant", "content": r["RESPONSE_TEXT"]})
        return messages
    except:
        return []


# -------------  SIDEBAR -------------
with st.sidebar:
    # Enhanced header
    st.markdown(
        f"""
        <div style="padding:16px; background:{theme['bg_secondary']}; border-radius:16px; 
                     margin-bottom:20px; border:1px solid {theme['border']};
                     box-shadow:0 2px 8px rgba(0,0,0,0.04);">
            <div style="display:flex; align-items:center; gap:10px; margin-bottom:12px;">
                <div style="width:32px; height:32px; background:{theme['message_user_bg']}; 
                            border-radius:8px; display:flex; align-items:center; justify-content:center;
                            box-shadow:0 2px 6px {theme['button_hover_shadow']};">
                    <span style="font-size:18px;">💬</span>
                </div>
                <h3 style="margin:0; color:{theme['text_primary']}; font-size:18px; font-weight:600;">
                    AIFAQ Pro
                </h3>
            </div>
        </div>
    """,
        unsafe_allow_html=True,
    )

    st.markdown("---")

    # Enhanced Knowledge Base Management header
    st.markdown(
        f"""
        <div style="display:flex; align-items:center; gap:10px; margin-bottom:16px;">
            <span style="font-size:28px;">📚</span>
            <h2 style="margin:0; color:{theme['text_primary']}; font-size:22px; font-weight:700;">
                Knowledge Base Management
            </h2>
        </div>
    """,
        unsafe_allow_html=True,
    )

    st.markdown("### 📤 Upload Documents")

    # File uploader outside the processing logic
    files = st.file_uploader(
        "Choose files",
        type=list(SUPPORTED_EXTENSIONS.keys()),
        accept_multiple_files=True,
        key="file_uploader",
    )

    # Upload button and processing
    if files and st.button(
        "⚡ Upload & Process", use_container_width=True, type="primary"
    ):
        progress_bar = st.progress(0)
        status_container = st.container()

        success_count = 0
        failed_files = []

        for idx, file in enumerate(files):
            with status_container:
                st.info(f"⏳ Processing {file.name}...")

            ext = file.name.split(".")[-1].lower()
            file_type = SUPPORTED_EXTENSIONS.get(ext, "TEXT")

            # Extract content - this now reads 100% of all file types
            content = extract_file_content(file, file_type)

            if content:
                # All documents are now private by default
                if store_document_fast(file.name, file_type, content, public=False):
                    with status_container:
                        st.success(f"✅ {file.name} uploaded successfully!")
                    success_count += 1

                    # CRITICAL: Force immediate rerun after each successful upload
                    st.session_state.kb_refresh = True
                    st.session_state.last_upload_time = time.time()
                    progress_bar.progress((idx + 1) / len(files))
                    time.sleep(0.3)  # Brief pause to show success message
                    st.rerun()  # Immediate refresh to update knowledge base
                else:
                    with status_container:
                        st.error(f"❌ {file.name} failed to upload")
                    failed_files.append(file.name)
            else:
                with status_container:
                    st.warning(f"⚠️ {file.name} - no content extracted")
                failed_files.append(file.name)

            progress_bar.progress((idx + 1) / len(files))

        # This code only runs if all files processed without triggering rerun
        progress_bar.empty()

        if success_count > 0:
            st.success(f"✨ All done! Successfully uploaded {success_count} file(s)!")
            if failed_files:
                st.warning(f"⚠️ Failed: {', '.join(failed_files)}")
        else:
            st.error("❌ No files were successfully processed.")
            if failed_files:
                st.error(f"Failed files: {', '.join(failed_files)}")

    st.markdown("---")
    st.markdown("### 📄 Knowledge Base")

    # Create placeholder for knowledge base that can be updated
    kb_container = st.container()

    with kb_container:
        # Show upload success indicator
        if st.session_state.kb_refresh:
            st.success("✅ Knowledge base updated!")
            st.session_state.kb_refresh = False

        # Force fresh database query - this ensures latest data
        df = get_user_docs()

        if df.empty:
            st.info("Your knowledge base is empty. Upload some files to get started!")
        else:
            st.caption(f"✨ Total documents: **{len(df)}**")

            # Display most recent uploads first
            for _, r in df.iterrows():
                type_icons = {
                    "PDF": "📄",
                    "PowerPoint": "📊",
                    "CSV": "📈",
                    "Excel": "📊",
                    "TEXT": "📝",
                }
                icon = type_icons.get(r["FILE_TYPE"], "📄")

                with st.expander(f"{icon} {r['FILENAME']}", expanded=False):
                    st.caption(f"**Type:** {r['FILE_TYPE']}")
                    st.caption(f"**Size:** {r['FILE_SIZE']:,} characters")
                    st.caption(
                        f"**Uploaded:** {r['UPLOADED_AT'].strftime('%Y-%m-%d %H:%M')}"
                    )

                    if st.button(
                        "🗑️ Delete", key=f"del_{r['DOC_ID']}", use_container_width=True
                    ):
                        if delete_document(r["DOC_ID"]):
                            st.success("Document deleted!")
                            time.sleep(0.5)
                            st.rerun()
                        else:
                            st.error("Failed to delete document")

    st.markdown("---")
    st.markdown("### 💬 Chat Sessions")
    if st.button("➕ New chat", use_container_width=True):
        st.session_state.current = str(uuid.uuid4())
        st.session_state.sessions[st.session_state.current] = {
            "title": "New chat",
            "messages": [],
        }
        st.rerun()

    try:
        chats = session.sql(
            f"SELECT DISTINCT SESSION_ID, LEFT(QUERY_TEXT,40) TITLE, MAX(QUERY_TIMESTAMP) TS "
            f"FROM {CHAT_TBL} "
            f"GROUP BY SESSION_ID, QUERY_TEXT ORDER BY TS DESC LIMIT 15"
        ).to_pandas()

        if not chats.empty:
            for _, c in chats.iterrows():
                btn_label = c["TITLE"] + "..." if len(c["TITLE"]) >= 40 else c["TITLE"]
                if st.button(
                    f"💬 {btn_label}",
                    key=f"chat_{c['SESSION_ID']}",
                    use_container_width=True,
                ):
                    st.session_state.current = c["SESSION_ID"]
                    # Load chat history from database
                    loaded_messages = load_chat_history(c["SESSION_ID"])
                    st.session_state.sessions[c["SESSION_ID"]] = {
                        "title": btn_label,
                        "messages": loaded_messages,
                    }
                    st.rerun()
    except Exception:
        pass

    if st.button("🗑️ Clear current chat", use_container_width=True):
        try:
            session.sql(
                f"DELETE FROM {CHAT_TBL} WHERE SESSION_ID='{sanitize(st.session_state.current)}'"
            ).collect()
            st.session_state.sessions[st.session_state.current]["messages"].clear()
            st.success("Chat history cleared!")
            st.rerun()
        except Exception as e:
            st.error(f"Error clearing chat: {e}")

# -------------  MAIN CHAT INTERFACE -------------
st.markdown(
    f"""
    <div style="margin-bottom:32px; padding:24px; background:{theme['bg_secondary']}; 
                border-radius:20px; border:1px solid {theme['border']};
                box-shadow:0 4px 12px rgba(0,0,0,0.04);">
        <div style="display:flex; align-items:center; gap:12px; margin-bottom:8px;">
            <div style="width:40px; height:40px; background:{theme['message_user_bg']}; 
                        border-radius:12px; display:flex; align-items:center; justify-content:center;
                        box-shadow:0 4px 8px {theme['button_hover_shadow']};">
                <span style="font-size:24px;">💬</span>
            </div>
            <h2 style="color:{theme['text_primary']};font-weight:700;margin:0;font-size:28px;">
                Chat with AIFAQ Pro
            </h2>
        </div>
        <p style="color:{theme['text_secondary']};margin:0;font-size:15px;padding-left:52px;">
            Private and secure AI-powered conversations with your enterprise data
        </p>
    </div>
""",
    unsafe_allow_html=True,
)

llm_selected = st.selectbox(
    "🧠 Choose Language Model",
    LLM_OPTIONS,
    index=LLM_OPTIONS.index(st.session_state.llm),
    key="llm_selector",
)
st.session_state.llm = llm_selected

st.markdown("---")

# Ensure current session exists
if st.session_state.current not in st.session_state.sessions:
    st.session_state.sessions[st.session_state.current] = {
        "title": "New chat",
        "messages": [],
    }

messages = st.session_state.sessions[st.session_state.current]["messages"]

# Display messages - Beautiful modern design with enhanced styling
chat_container = st.container()
with chat_container:
    if not messages:
        st.markdown(
            f"""
            <div style="text-align:center; padding:80px 20px;">
                <div style="display:inline-block; padding:20px; border-radius:20px; 
                            background:{theme['bg_secondary']}; 
                            box-shadow:0 8px 24px rgba(0,0,0,0.06);
                            border:1px solid {theme['border']};">
                    <h2 style="color:{theme['primary']}; font-weight:600; margin:0 0 12px 0;">
                        👋 Welcome to AIFAQ Pro AI Assistant
                    </h2>
                    <p style="color:{theme['text_secondary']}; font-size:16px; margin:0; line-height:1.6;">
                        Upload to your knowledge base and start asking questions to unlock insights<br>
                        <span style="font-size:14px; color:{theme['text_secondary']}; opacity:0.8;">
                        Powered by advanced AI • Fast • Secure
                        </span>
                    </p>
                </div>
            </div>
        """,
            unsafe_allow_html=True,
        )
    else:
        for idx, msg in enumerate(messages):
            if msg["role"] == "user":
                st.markdown(
                    f'<div style="display:flex;justify-content:flex-end;margin:20px 0;animation:slideIn 0.3s ease-out;">'
                    f'<div style="background:{theme["message_user_bg"]};'
                    f'color:{theme["message_user_text"]};'
                    f"border-radius:20px 20px 4px 20px;"
                    f"padding:16px 20px;"
                    f"max-width:75%;"
                    f'box-shadow:0 4px 12px {theme["button_hover_shadow"]}, 0 2px 4px rgba(0,0,0,0.08);'
                    f"font-size:15px;"
                    f"line-height:1.6;"
                    f"font-weight:500;"
                    f"border:1px solid rgba(255,255,255,0.2);"
                    f"backdrop-filter:blur(10px);"
                    f'">'
                    f'{msg["content"]}</div></div>',
                    unsafe_allow_html=True,
                )
            else:
                st.markdown(
                    f'<div style="display:flex;justify-content:flex-start;margin:20px 0;animation:slideIn 0.3s ease-out;">'
                    f'<div style="background:{theme["message_ai_bg"]};'
                    f'color:{theme["message_ai_text"]};'
                    f"border-radius:20px 20px 20px 4px;"
                    f"padding:16px 20px;"
                    f"max-width:75%;"
                    f"box-shadow:0 2px 8px rgba(0,0,0,0.06), 0 4px 16px rgba(0,0,0,0.04);"
                    f"font-size:15px;"
                    f"line-height:1.7;"
                    f'border:1px solid {theme["message_ai_border"]};'
                    f'position:relative;">'
                    f'<div style="position:absolute; top:-6px; left:-6px; width:12px; height:12px; '
                    f'background:{theme["accent_highlight"]}; border-radius:50%; opacity:0.6;"></div>'
                    f'{msg["content"]}</div></div>',
                    unsafe_allow_html=True,
                )

        # Add CSS animation
        st.markdown(
            """
            <style>
            @keyframes slideIn {
                from {
                    opacity: 0;
                    transform: translateY(10px);
                }
                to {
                    opacity: 1;
                    transform: translateY(0);
                }
            }
            </style>
        """,
            unsafe_allow_html=True,
        )

# Input form - Beautiful modern design with enhanced styling
with st.form("chat_form", clear_on_submit=True):
    # Enhanced input styling
    st.markdown(
        f"""
        <style>
        .stTextArea textarea {{
            border-radius: 20px !important;
            border: 2px solid {theme['input_border']} !important;
            padding: 18px 24px !important;
            font-size: 15px !important;
            background: {theme['input_bg']} !important;
            box-shadow: 0 2px 8px rgba(0,0,0,0.04) !important;
            transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1) !important;
        }}
        .stTextArea textarea:focus {{
            border-color: {theme['input_focus']} !important;
            box-shadow: 0 0 0 4px rgba(139, 92, 246, 0.15), 
                        0 8px 16px rgba(0,0,0,0.08) !important;
            transform: translateY(-2px);
        }}
        .stTextArea textarea::placeholder {{
            color: {theme['text_secondary']} !important;
            opacity: 0.6;
        }}
        </style>
    """,
        unsafe_allow_html=True,
    )

    prompt = st.text_area(
        "message",
        placeholder="🔍  Type your question or search your knowledge base...",
        height=90,
        label_visibility="collapsed",
        key="chat_input",
    )

    # Action buttons with enhanced design
    col1, col2, col3 = st.columns([6, 1.5, 1])
    with col1:
        st.markdown(
            f"""
            <div style="display:flex; align-items:center; gap:8px; padding:8px 0;">
                <span style="color:{theme['text_secondary']}; font-size:14px;">
                    💡 <strong>Tip:</strong> Be specific and detailed for best results
                </span>
            </div>
        """,
            unsafe_allow_html=True,
        )
    with col2:
        send = st.form_submit_button(
            "✨ Send", use_container_width=True, type="primary"
        )
    with col3:
        clear = st.form_submit_button("Clear", use_container_width=True)

    if clear:
        st.session_state.sessions[st.session_state.current]["messages"].clear()
        st.rerun()

    if send and prompt.strip():
        messages.append({"role": "user", "content": prompt.strip()})

        with st.spinner("🧠 Thinking..."):
            try:
                # Search for relevant document chunks
                rows = session.sql(
                    f"""
                    WITH qv AS (
                        SELECT SNOWFLAKE.CORTEX.EMBED_TEXT_768('snowflake-arctic-embed-m',
                            LEFT('{sanitize(prompt.strip())}', 10000)) q_vec
                    ),
                    sim AS (
                        SELECT c.CHUNK_TEXT, d.FILENAME, d.FILE_TYPE,
                                VECTOR_COSINE_SIMILARITY(e.EMBEDDING, qv.q_vec) AS SIMILARITY_SCORE
                        FROM {CHUNKS_TBL} c
                        JOIN {EMB_TBL} e ON c.CHUNK_ID = e.CHUNK_ID
                        JOIN {DOCS_TBL} d ON c.DOC_ID = d.DOC_ID
                        CROSS JOIN qv
                        WHERE VECTOR_COSINE_SIMILARITY(e.EMBEDDING, qv.q_vec) > 0.3
                        ORDER BY SIMILARITY_SCORE DESC
                        LIMIT 10 
                    )
                    SELECT CHUNK_TEXT, FILENAME, FILE_TYPE, SIMILARITY_SCORE FROM sim
                """
                ).to_pandas()

                if rows.empty:
                    answer = "I couldn't find any relevant information in your knowledge base to answer your question. Please upload some files first or try rephrasing your question."
                else:
                    # Create context from relevant chunks
                    context_parts = []
                    for _, row in rows.iterrows():
                        context_parts.append(
                            f"[From {row['FILENAME']} ({row['FILE_TYPE']})]\n{row['CHUNK_TEXT']}"
                        )
                    context = "\n\n---\n\n".join(context_parts)

                    # Generate response using LLM with source tracking
                    prompt_llm = (
                        f"You are a helpful assistant that answers questions based on the provided knowledge base context. "
                        f"Answer the user's question using only the information from the provided knowledge base. "
                        f"If the context contains numerical data, tables, or statistics, include them in your answer. "
                        f"Be concise but comprehensive.\n\n"
                        f"Context:\n{context}\n\n"
                        f"Question: {prompt.strip()}\n\n"
                        f"Answer:"
                    )

                    ans_rows = session.sql(
                        f"SELECT SNOWFLAKE.CORTEX.COMPLETE('{sanitize(st.session_state.llm)}', "
                        f"'{sanitize(prompt_llm)}') as ans"
                    ).collect()
                    answer = (
                        str(ans_rows[0][0])
                        if ans_rows
                        else "The LLM service didn't return a response."
                    )

                    # Cite only the PRIMARY source (highest similarity score)
                    # Group by filename and get the highest similarity for each document
                    source_scores = (
                        rows.groupby("FILENAME")["SIMILARITY_SCORE"]
                        .max()
                        .sort_values(ascending=False)
                    )

                    # Get top source(s) - only cite documents with similarity > 0.5 (strong match)
                    top_sources = []
                    for filename, score in source_scores.items():
                        if score > 0.5:  # Strong relevance threshold
                            top_sources.append(filename)
                            if len(top_sources) >= 2:  # Max 2 sources
                                break

                    # If no strong matches, just use the top result
                    if not top_sources:
                        top_sources = [source_scores.index[0]]

                    # Add source attribution - only the primary source(s) used
                    if len(top_sources) == 1:
                        sources_md = f"\n\n---\n**Source:** {top_sources[0]}"
                    else:
                        sources_list_md = "\n- ".join(top_sources)
                        sources_md = (
                            f"\n\n---\n" f"**Sources:**\n" f"- {sources_list_md}"
                        )
                    answer += sources_md

            except Exception as e:
                answer = f"❌ Error processing your request: {str(e)}"

        messages.append({"role": "assistant", "content": answer})

        # Persist chat history
        try:
            session.write_pandas(
                pd.DataFrame(
                    [
                        {
                            "CHAT_ID": str(uuid.uuid4()),
                            "SESSION_ID": st.session_state.current,
                            "QUERY_TEXT": prompt.strip(),
                            "RESPONSE_TEXT": answer,
                            "QUERY_TIMESTAMP": dt.datetime.utcnow().strftime(
                                "%Y-%m-%d %H:%M:%S"
                            ),
                        }
                    ]
                ),
                table_name=CHAT_TBL,
                auto_create_table=False,
                overwrite=False,
            )
        except Exception:
            pass  # Chat history persistence is optional

        st.rerun()